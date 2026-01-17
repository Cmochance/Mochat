"""
Cloud Run PPT 生成服务 - 基于模板版本（简化版）
支持幻灯片类型: title, toc, section, content, two-column, thank-you
"""
import functions_framework
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import boto3
import os
import io
import json
import uuid
from datetime import datetime

# 模板配置
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), 'templates', 'sci.pptx')

# 模板幻灯片索引 (0-based)
TEMPLATE_SLIDES = {
    'title': 0,       # 封面页
    'toc': 1,         # 目录页
    'section': 2,     # 章节分隔页
    'content': 3,     # 内容页
    'two-column': 4,  # 双栏页
    'thank-you': 5,   # 结束页
}


def hex_to_rgb(hex_color):
    """将十六进制颜色转换为 RGBColor"""
    if not hex_color:
        return RGBColor(0, 0, 0)
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def set_font(run, size=None, color=None, bold=None, italic=None, font_name='Microsoft YaHei'):
    """统一设置字体属性，包括中文字体"""
    if size is not None:
        run.font.size = Pt(size)
    if color is not None:
        run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    run.font.name = font_name
    
    # 设置东亚字体
    r = run._r
    rPr = r.find(qn('a:rPr'))
    if rPr is None:
        rPr = etree.SubElement(r, qn('a:rPr'))
        r.insert(0, rPr)
    
    for tag in ['a:ea', 'a:latin', 'a:cs']:
        elem = rPr.find(qn(tag))
        if elem is None:
            elem = etree.SubElement(rPr, qn(tag))
        elem.set('typeface', font_name)


def set_shape_text(shape, text, font_name='Microsoft YaHei', color=None):
    """设置形状中的文本 - 完全清空并重写"""
    if not hasattr(shape, 'text_frame'):
        return
    tf = shape.text_frame
    
    # 清空所有段落的所有文本
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ''
    
    # 如果有文本要设置，在第一个段落添加
    if text:
        if tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].text = text
            set_font(tf.paragraphs[0].runs[0], font_name=font_name, color=color)
        else:
            run = tf.paragraphs[0].add_run()
            run.text = text
            set_font(run, font_name=font_name, color=color)


def add_text_box(slide, left, top, width, height, text, size, color, bold=False, italic=False, align=None, font_name='Microsoft YaHei'):
    """添加文本框"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold, italic, font_name=font_name)
    return box


# ============ 渲染函数 ============

def render_title_slide(slide, slide_data, theme):
    """渲染封面页"""
    font_name = theme.get('fontFamily', 'Microsoft YaHei')
    
    title = slide_data.get('title', '演示文稿')
    subtitle = slide_data.get('subtitle', '')
    date_str = datetime.now().strftime('%Y年%m月%d日')
    
    title_shape = None
    subtitle_shape = None
    
    for shape in slide.shapes:
        if not hasattr(shape, 'text'):
            continue
        text = shape.text
        
        if '宣讲会' in text or '材料与化学学院宣讲会' in text:
            set_shape_text(shape, title, font_name)
            title_shape = shape
        elif 'CHINA' in text or 'GEOSCIENCES' in text:
            set_shape_text(shape, subtitle, font_name)
            subtitle_shape = shape
        elif 'XXXX年' in text:
            set_shape_text(shape, date_str, font_name)
        elif '材料与化学学院' in text and '宣讲会' not in text:
            set_shape_text(shape, '', font_name)
        elif '言简意赅' in text or '阐述观点' in text:
            set_shape_text(shape, '', font_name)
    
    # 调整副标题位置，在主标题下方留出0.5行间距
    if title_shape and subtitle_shape:
        # 计算主标题底部位置 + 0.3英寸间距
        title_bottom = title_shape.top + title_shape.height
        subtitle_shape.top = title_bottom + Inches(0.3)


def render_toc_slide(slide, slide_data, theme):
    """渲染目录页"""
    font_name = theme.get('fontFamily', 'Microsoft YaHei')
    
    toc_title = slide_data.get('title', '目录')
    items = slide_data.get('items', [])
    
    group_index = 0
    for shape in slide.shapes:
        if hasattr(shape, 'text') and '目录' in shape.text:
            set_shape_text(shape, toc_title, font_name)
        
        if shape.shape_type == 6:  # GROUP
            if group_index < len(items):
                item = items[group_index]
                # 支持新格式（对象）和旧格式（字符串）
                if isinstance(item, dict):
                    item_title = item.get('title', '')
                    item_summary = item.get('summary', '')
                else:
                    item_title = str(item)
                    item_summary = ''
                
                for sub_shape in shape.shapes:
                    if hasattr(sub_shape, 'text'):
                        text = sub_shape.text
                        if text in ('01', '02', '03', '04'):
                            set_shape_text(sub_shape, f"{group_index + 1:02d}", font_name)
                        elif '添加标题' in text:
                            # 设置标题和摘要两行显示
                            if item_summary:
                                # 清空原有内容，添加两行
                                tf = sub_shape.text_frame
                                for para in tf.paragraphs:
                                    for run in para.runs:
                                        run.text = ''
                                
                                # 第一行：标题（较大字体）
                                if tf.paragraphs[0].runs:
                                    tf.paragraphs[0].runs[0].text = item_title
                                    set_font(tf.paragraphs[0].runs[0], size=16, font_name=font_name)
                                else:
                                    run = tf.paragraphs[0].add_run()
                                    run.text = item_title
                                    set_font(run, size=16, font_name=font_name)
                                
                                # 第二行：摘要（较小字体）
                                p = tf.paragraphs[0]._p
                                from pptx.oxml.ns import qn
                                new_p = etree.SubElement(p.getparent(), qn('a:p'))
                                p.getparent().append(new_p)
                                
                                from pptx.text.text import _Paragraph
                                new_para = _Paragraph(new_p, tf.paragraphs[0]._parent)
                                summary_run = new_para.add_run()
                                summary_run.text = item_summary
                                set_font(summary_run, size=12, font_name=font_name, 
                                        color=hex_to_rgb('#666666'))
                            else:
                                set_shape_text(sub_shape, item_title, font_name)
                        elif '点击此处' in text or '言简意赅' in text:
                            set_shape_text(sub_shape, '', font_name)
            else:
                for sub_shape in shape.shapes:
                    if hasattr(sub_shape, 'text'):
                        set_shape_text(sub_shape, '', font_name)
            group_index += 1
    
    # 超过4个章节的处理（保留兼容性）
    if len(items) > 4:
        y_pos = 1.9
        for i, item in enumerate(items[4:6], start=5):
            if isinstance(item, dict):
                display_text = f"{i:02d}  {item.get('title', '')}"
            else:
                display_text = f"{i:02d}  {item}"
            add_text_box(slide, 7.8, y_pos + (i-5) * 0.6, 5, 0.5,
                        display_text, 14, hex_to_rgb(theme.get('textColor', '#333333')), 
                        font_name=font_name)


def render_section_slide(slide, slide_data, theme):
    """渲染章节分隔页"""
    font_name = theme.get('fontFamily', 'Microsoft YaHei')
    title = slide_data.get('title', '章节')
    
    shapes_to_remove = []
    
    for shape in slide.shapes:
        if not hasattr(shape, 'text'):
            continue
        text = shape.text
        
        # 主标题 - "编辑标题"
        if '编辑标题' in text:
            set_shape_text(shape, title, font_name)
        # 副标题和其他文本 - 删除
        elif '副标题' in text or '单击此处' in text or '添加文本' in text:
            shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def render_content_slide(slide, slide_data, theme):
    """渲染内容页"""
    font_name = theme.get('fontFamily', 'Microsoft YaHei')
    text_color = hex_to_rgb(theme.get('textColor', '#333333'))
    
    title = slide_data.get('title', '')
    content_items = slide_data.get('content', [])
    
    # 需要删除的形状
    shapes_to_remove = []
    
    for shape in slide.shapes:
        if not hasattr(shape, 'text'):
            continue
        text = shape.text
        
        if '单击此处添加标题' in text:
            set_shape_text(shape, title, font_name)
        elif '单击此处添加正文' in text:
            # 这是内容占位符，需要删除
            shapes_to_remove.append(shape)
        elif '点击此处添加正文' in text or '言简意赅' in text:
            shapes_to_remove.append(shape)
    
    # 删除占位符形状
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    
    y_pos = 1.7
    for item in content_items:
        if isinstance(item, str):
            add_text_box(slide, 0.7, y_pos, 12, 0.6, item, 18, text_color, font_name=font_name)
            y_pos += 0.7
        elif isinstance(item, dict):
            item_type = item.get('type', 'text')
            if item_type == 'text':
                add_text_box(slide, 0.7, y_pos, 12, 0.6, item.get('value', ''), 18, text_color, font_name=font_name)
                y_pos += 0.7
            elif item_type in ('bullet', 'numbered'):
                for i, txt in enumerate(item.get('items', [])):
                    prefix = f"{i+1}. " if item_type == 'numbered' else "• "
                    add_text_box(slide, 0.9, y_pos, 11.8, 0.5, prefix + txt, 16, text_color, font_name=font_name)
                    y_pos += 0.55
                y_pos += 0.15


def render_two_column_slide(slide, slide_data, theme):
    """渲染双栏页"""
    font_name = theme.get('fontFamily', 'Microsoft YaHei')
    text_color = hex_to_rgb(theme.get('textColor', '#333333'))
    
    title = slide_data.get('title', '')
    left_items = slide_data.get('left', [])
    right_items = slide_data.get('right', [])
    
    shapes_to_remove = []
    
    for shape in slide.shapes:
        if not hasattr(shape, 'text'):
            continue
        text = shape.text
        
        if '单击此处添加标题' in text:
            set_shape_text(shape, title, font_name)
        elif '点击此处添加正文' in text or '言简意赅' in text:
            shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    
    y_pos = 2.1
    for item in left_items:
        if isinstance(item, dict):
            item_type = item.get('type', 'text')
            if item_type == 'text':
                add_text_box(slide, 1.4, y_pos, 5.0, 0.5, item.get('value', ''), 16, text_color, font_name=font_name)
                y_pos += 0.6
            elif item_type in ('bullet', 'numbered'):
                for i, txt in enumerate(item.get('items', [])):
                    prefix = f"{i+1}. " if item_type == 'numbered' else "• "
                    add_text_box(slide, 1.5, y_pos, 4.9, 0.45, prefix + txt, 14, text_color, font_name=font_name)
                    y_pos += 0.5
    
    y_pos = 2.1
    for item in right_items:
        if isinstance(item, dict):
            item_type = item.get('type', 'text')
            if item_type == 'text':
                add_text_box(slide, 6.9, y_pos, 5.0, 0.5, item.get('value', ''), 16, text_color, font_name=font_name)
                y_pos += 0.6
            elif item_type in ('bullet', 'numbered'):
                for i, txt in enumerate(item.get('items', [])):
                    prefix = f"{i+1}. " if item_type == 'numbered' else "• "
                    add_text_box(slide, 7.0, y_pos, 4.9, 0.45, prefix + txt, 14, text_color, font_name=font_name)
                    y_pos += 0.5


def render_quote_slide(slide, slide_data, theme):
    """渲染引用页"""
    font_name = theme.get('fontFamily', 'Microsoft YaHei')
    primary_color = hex_to_rgb(theme.get('primaryColor', '#1a73e8'))
    text_color = hex_to_rgb(theme.get('textColor', '#333333'))
    
    quote = slide_data.get('quote', '')
    author = slide_data.get('author', '')
    
    for shape in slide.shapes:
        if not hasattr(shape, 'text'):
            continue
        text = shape.text
        if '单击此处添加标题' in text or '单击此处添加正文' in text or '点击此处添加正文' in text or '言简意赅' in text:
            set_shape_text(shape, '', font_name)
    
    add_text_box(slide, 1.5, 2.5, 10.3, 2, f'"{quote}"', 28, primary_color, 
                italic=True, align=PP_ALIGN.CENTER, font_name=font_name)
    if author:
        add_text_box(slide, 1.5, 5.0, 10.3, 0.5, f"—— {author}", 18, text_color,
                    align=PP_ALIGN.RIGHT, font_name=font_name)


def render_thank_you_slide(slide, slide_data, theme):
    """渲染结束页"""
    font_name = theme.get('fontFamily', 'Microsoft YaHei')
    
    title = slide_data.get('title', 'THANKS')
    subtitle = slide_data.get('subtitle', '')
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text = shape.text
            if 'THANKS' in text:
                set_shape_text(shape, title, font_name)
            elif '添加文本' in text:
                set_shape_text(shape, subtitle, font_name)


def generate_ppt(ppt_data):
    """生成 PPT - 直接在模板上修改，不删除幻灯片"""
    # 主题配置
    default_theme = {
        'primaryColor': '#1a73e8',
        'secondaryColor': '#34a853',
        'backgroundColor': '#ffffff',
        'textColor': '#333333',
        'fontFamily': 'Microsoft YaHei'
    }
    theme = {**default_theme, **ppt_data.get('theme', {})}
    
    slides_data = ppt_data.get('slides', [])
    
    # 加载模板
    prs = Presentation(TEMPLATE_PATH)
    
    # 模板有6页: 0=title, 1=toc, 2=section, 3=content, 4=two-column, 5=thank-you
    # 记录每种类型使用了哪个模板幻灯片
    type_to_slide = {}
    
    for sd in slides_data:
        slide_type = sd.get('type', 'content')
        if slide_type == 'quote':
            slide_type = 'content'
        
        template_idx = TEMPLATE_SLIDES.get(slide_type, TEMPLATE_SLIDES['content'])
        
        if slide_type not in type_to_slide:
            # 第一次使用这种类型，使用模板中的幻灯片
            slide = prs.slides[template_idx]
            type_to_slide[slide_type] = template_idx
        else:
            # 已经使用过，需要添加新幻灯片（使用布局，不复制内容）
            template_slide = prs.slides[template_idx]
            slide = prs.slides.add_slide(template_slide.slide_layout)
        
        # 渲染内容
        if slide_type == 'title':
            render_title_slide(slide, sd, theme)
        elif slide_type == 'toc':
            render_toc_slide(slide, sd, theme)
        elif slide_type == 'section':
            render_section_slide(slide, sd, theme)
        elif slide_type == 'content':
            if sd.get('type') == 'quote':
                render_quote_slide(slide, sd, theme)
            else:
                render_content_slide(slide, sd, theme)
        elif slide_type == 'two-column':
            render_two_column_slide(slide, sd, theme)
        elif slide_type == 'thank-you':
            render_thank_you_slide(slide, sd, theme)
    
    # 不删除未使用的幻灯片，保持文件完整性
    # 用户可以手动删除不需要的幻灯片
    
    return prs


@functions_framework.http
def pptgen(request):
    """Cloud Run PPT 生成服务入口"""
    # 调试端点
    if request.args.get('debug') == 'template':
        info = {
            "template_path": TEMPLATE_PATH,
            "template_exists": os.path.exists(TEMPLATE_PATH),
            "template_size": os.path.getsize(TEMPLATE_PATH) if os.path.exists(TEMPLATE_PATH) else 0,
            "cwd": os.getcwd(),
            "dirname": os.path.dirname(__file__),
            "files_in_templates": os.listdir(os.path.join(os.path.dirname(__file__), 'templates')) if os.path.exists(os.path.join(os.path.dirname(__file__), 'templates')) else []
        }
        return (json.dumps(info), 200, {'Content-Type': 'application/json'})
    
    if request.args.get('debug') == 'versions':
        import pptx
        import lxml
        info = {
            "python_pptx": pptx.__version__,
            "lxml": lxml.__version__,
            "boto3": boto3.__version__
        }
        return (json.dumps(info), 200, {'Content-Type': 'application/json'})
    
    # 1. CORS
    if request.method == 'OPTIONS':
        headers = {
            'Access-Control-Allow-Origin': '*',
            'Access-Control-Allow-Methods': 'POST',
            'Access-Control-Allow-Headers': 'Content-Type, Authorization, X-Auth-Secret',
            'Access-Control-Max-Age': '3600'
        }
        return ('', 204, headers)

    headers = {
        'Access-Control-Allow-Origin': '*',
        'Content-Type': 'application/json'
    }

    # 2. 鉴权
    auth_secret = request.headers.get("X-Auth-Secret") or request.headers.get("Authorization")
    env_token = os.environ.get("AUTH_TOKEN") or os.environ.get("PPTGEN_CLOUDRUN_SECRET", "")
    if env_token and auth_secret != env_token:
        return (json.dumps({"error": "Unauthorized"}), 401, headers)

    # 3. 解析请求
    try:
        request_json = request.get_json(silent=True)
        if not request_json:
            return (json.dumps({"error": "Invalid JSON"}), 400, headers)
        
        ppt_data = request_json.get("ppt_data", request_json)
        filename = request_json.get("filename", ppt_data.get("title", "presentation"))
        user_id = request_json.get("user_id", "anonymous")
        
        if not ppt_data or 'slides' not in ppt_data:
            return (json.dumps({"error": "Missing slides data"}), 400, headers)
            
    except Exception as e:
        return (json.dumps({"error": str(e)}), 400, headers)

    # 4. 生成 PPT
    try:
        prs = generate_ppt(ppt_data)
        
        # 保存到临时文件而不是 BytesIO
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp_path = tmp.name
        
        prs.save(tmp_path)
        
        # 读取文件内容
        with open(tmp_path, 'rb') as f:
            ppt_bytes = f.read()
        
        # 删除临时文件
        os.unlink(tmp_path)
        
    except Exception as e:
        import traceback
        return (json.dumps({"error": f"PPT Generation Error: {str(e)}", "trace": traceback.format_exc()}), 500, headers)

    # 5. 上传到 R2
    try:
        from botocore.config import Config
        
        s3 = boto3.client(
            's3',
            endpoint_url=f"https://{os.environ.get('R2_ACCOUNT_ID')}.r2.cloudflarestorage.com",
            aws_access_key_id=os.environ.get('R2_ACCESS_KEY_ID'),
            aws_secret_access_key=os.environ.get('R2_SECRET_ACCESS_KEY'),
            config=Config(signature_version='s3v4')
        )

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        unique_id = str(uuid.uuid4())[:8]
        safe_filename = "".join(c for c in filename if c.isalnum() or c in (' ', '-', '_') or ord(c) > 127)
        safe_filename = safe_filename[:50] or "presentation"
        
        r2_key = f"ppt/{user_id}/{timestamp}_{unique_id}_{safe_filename}.pptx"
        
        # 使用 put_object 上传
        s3.put_object(
            Bucket=os.environ.get('R2_BUCKET_NAME'),
            Key=r2_key,
            Body=ppt_bytes,
            ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            ContentDisposition=f'attachment; filename="{safe_filename}.pptx"'
        )
        
        # 生成 presigned URL 直接从 R2 下载，绕过 Cloudflare CDN 转换
        # URL 有效期 7 天
        presigned_url = s3.generate_presigned_url(
            'get_object',
            Params={
                'Bucket': os.environ.get('R2_BUCKET_NAME'),
                'Key': r2_key,
                'ResponseContentDisposition': f'attachment; filename="{safe_filename}.pptx"'
            },
            ExpiresIn=7 * 24 * 3600  # 7 天
        )
        
        return (json.dumps({
            "status": "success",
            "url": presigned_url,
            "title": ppt_data.get('title', '演示文稿')
        }), 200, headers)

    except Exception as e:
        return (json.dumps({"error": f"Upload R2 Error: {str(e)}"}), 500, headers)
